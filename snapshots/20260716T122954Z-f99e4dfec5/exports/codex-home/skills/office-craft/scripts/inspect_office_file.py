"""Inspect common office files and print a compact JSON summary."""

from __future__ import annotations

import json
import sys
from pathlib import Path


def inspect_docx(path: Path) -> dict:
    from docx import Document

    doc = Document(str(path))
    headings = []
    for paragraph in doc.paragraphs:
        style = paragraph.style.name if paragraph.style is not None else ""
        text = paragraph.text.strip()
        if text and style.lower().startswith("heading"):
            headings.append(text)
    return {
        "type": "docx",
        "paragraphs": len(doc.paragraphs),
        "tables": len(doc.tables),
        "headings": headings[:20],
    }


def inspect_xlsx(path: Path) -> dict:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=False, data_only=False)
    sheets = []
    for ws in wb.worksheets:
        formulas = 0
        nonempty = 0
        for row in ws.iter_rows():
            for cell in row:
                if cell.value is not None:
                    nonempty += 1
                    if isinstance(cell.value, str) and cell.value.startswith("="):
                        formulas += 1
        sheets.append(
            {
                "name": ws.title,
                "max_row": ws.max_row,
                "max_column": ws.max_column,
                "nonempty_cells": nonempty,
                "formulas": formulas,
            }
        )
    return {"type": "xlsx", "sheets": sheets}


def inspect_csv(path: Path) -> dict:
    import pandas as pd

    df = pd.read_csv(path)
    return {
        "type": "csv",
        "rows": int(df.shape[0]),
        "columns": int(df.shape[1]),
        "column_names": list(map(str, df.columns[:30])),
    }


def inspect_pdf(path: Path) -> dict:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    info = {"type": "pdf", "pages": len(reader.pages), "encrypted": reader.is_encrypted}
    if not reader.is_encrypted:
        sample = ""
        try:
            sample = (reader.pages[0].extract_text() or "")[:500] if reader.pages else ""
        except Exception as exc:
            sample = f"extract_error:{type(exc).__name__}:{exc}"
        info["first_page_text_sample"] = sample
    return info


def inspect_pptx(path: Path) -> dict:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = []
    for index, slide in enumerate(prs.slides, 1):
        texts = []
        pictures = 0
        tables = 0
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                tables += 1
            if getattr(shape, "shape_type", None) == 13:
                pictures += 1
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip().replace("\n", " | ")[:160])
        slides.append(
            {
                "index": index,
                "text_blocks": len(texts),
                "pictures": pictures,
                "tables": tables,
                "sample": texts[:3],
            }
        )
    return {"type": "pptx", "slides": len(prs.slides), "slide_summaries": slides}


def main(argv: list[str]) -> int:
    if len(argv) != 2:
        print("usage: inspect_office_file.py <file>", file=sys.stderr)
        return 2
    path = Path(argv[1])
    if not path.exists():
        print(json.dumps({"error": "not_found", "path": str(path)}, ensure_ascii=False))
        return 1
    suffix = path.suffix.lower()
    if suffix == ".docx":
        result = inspect_docx(path)
    elif suffix == ".xlsx":
        result = inspect_xlsx(path)
    elif suffix == ".csv":
        result = inspect_csv(path)
    elif suffix == ".pdf":
        result = inspect_pdf(path)
    elif suffix == ".pptx":
        result = inspect_pptx(path)
    else:
        result = {"type": "unsupported", "suffix": suffix}
    result["path"] = str(path)
    print(json.dumps(result, ensure_ascii=False, indent=2))
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
