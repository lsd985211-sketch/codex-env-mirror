"""Small, dependency-aware file analyzer for mobile bridge attachments.

This module is intentionally read-only. Editing user files belongs in a
separate explicit workflow with backups and confirmation.
"""

from __future__ import annotations

import csv
import hashlib
import json
import zipfile
from pathlib import Path
from typing import Any


MAX_PREVIEW_CHARS = 1200
MAX_TEXT_READ_BYTES = 2 * 1024 * 1024
MAX_ARCHIVE_MEMBERS = 30

TEXT_SUFFIXES = {
    ".txt",
    ".md",
    ".log",
    ".json",
    ".xml",
    ".yml",
    ".yaml",
    ".toml",
    ".ini",
    ".cfg",
    ".properties",
    ".csv",
}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tif", ".tiff"}
AUDIO_SUFFIXES = {".wav", ".mp3", ".m4a", ".aac", ".ogg", ".opus", ".amr", ".silk", ".aud", ".flac", ".wma"}
ZIP_SUFFIXES = {".zip", ".jar", ".mcpack", ".mcaddon", ".docx", ".xlsx", ".xlsm", ".pptx"}


def _truncate(text: str, limit: int = MAX_PREVIEW_CHARS) -> str:
    return text[:limit] + ("...<truncated>" if len(text) > limit else "")


def _json_preview(value: Any, limit: int = MAX_PREVIEW_CHARS) -> str:
    return _truncate(json.dumps(value, ensure_ascii=False), limit)


def _sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def _detect_encoding(path: Path) -> str:
    raw = path.read_bytes()[:MAX_TEXT_READ_BYTES]
    if raw.startswith(b"\xef\xbb\xbf"):
        return "utf-8-sig"
    try:
        raw.decode("utf-8")
        return "utf-8"
    except UnicodeDecodeError:
        pass
    try:
        from charset_normalizer import from_bytes  # type: ignore

        best = from_bytes(raw).best()
        if best and best.encoding:
            return str(best.encoding)
    except Exception:
        pass
    return "utf-8"


def _read_text_preview(path: Path) -> dict[str, Any]:
    encoding = _detect_encoding(path)
    try:
        raw = path.read_bytes()[:MAX_TEXT_READ_BYTES]
        text = raw.decode(encoding, errors="replace").replace("\x00", "")
    except Exception as exc:
        return {"kind": "text", "ok": False, "preview": f"preview_error={exc}"}
    return {
        "kind": "text",
        "ok": True,
        "encoding": encoding,
        "preview": _truncate(text),
        "truncated": path.stat().st_size > len(raw) or len(text) > MAX_PREVIEW_CHARS,
    }


def _csv_preview(path: Path) -> dict[str, Any]:
    encoding = _detect_encoding(path)
    try:
        with path.open("r", encoding=encoding, errors="replace", newline="") as fh:
            rows = []
            for index, row in enumerate(csv.reader(fh)):
                if index >= 8:
                    break
                rows.append(row[:8])
    except Exception as exc:
        return {"kind": "csv", "ok": False, "preview": f"preview_error={exc}"}
    return {"kind": "csv", "ok": True, "encoding": encoding, "preview": _json_preview(rows)}


def _xlsx_preview(path: Path) -> dict[str, Any]:
    try:
        import openpyxl  # type: ignore
    except Exception:
        return {"kind": "spreadsheet", "ok": False, "preview": "preview_unavailable=openpyxl is not installed"}
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets[:3]:
            rows = []
            for index, row in enumerate(sheet.iter_rows(values_only=True)):
                if index >= 8:
                    break
                rows.append([cell for cell in row[:8]])
            parts.append({"sheet": sheet.title, "rows": rows})
        wb.close()
        return {"kind": "spreadsheet", "ok": True, "preview": _json_preview(parts)}
    except Exception as exc:
        return {"kind": "spreadsheet", "ok": False, "preview": f"preview_error={exc}"}


def _xls_preview(path: Path) -> dict[str, Any]:
    try:
        import xlrd  # type: ignore
    except Exception:
        return {"kind": "spreadsheet", "ok": False, "preview": "preview_unavailable=xlrd is not installed"}
    try:
        book = xlrd.open_workbook(str(path))
        parts = []
        for sheet in book.sheets()[:3]:
            rows = []
            for row_index in range(min(sheet.nrows, 8)):
                rows.append(sheet.row_values(row_index, 0, min(sheet.ncols, 8)))
            parts.append({"sheet": sheet.name, "rows": rows})
        return {"kind": "spreadsheet", "ok": True, "preview": _json_preview(parts)}
    except Exception as exc:
        return {"kind": "spreadsheet", "ok": False, "preview": f"preview_error={exc}"}


def _docx_preview(path: Path) -> dict[str, Any]:
    try:
        import docx  # type: ignore
    except Exception:
        return {"kind": "document", "ok": False, "preview": "preview_unavailable=python-docx is not installed"}
    try:
        doc = docx.Document(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return {"kind": "document", "ok": True, "preview": _truncate("\n".join(paragraphs[:30]))}
    except Exception as exc:
        return {"kind": "document", "ok": False, "preview": f"preview_error={exc}"}


def _pptx_preview(path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return {"kind": "presentation", "ok": False, "preview": "preview_unavailable=python-pptx is not installed"}
    try:
        prs = Presentation(path)
        slides = []
        for slide_index, slide in enumerate(prs.slides, start=1):
            if slide_index > 5:
                break
            texts = []
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if isinstance(text, str) and text.strip():
                    texts.append(text.strip())
            slides.append({"slide": slide_index, "text": texts[:8]})
        return {"kind": "presentation", "ok": True, "preview": _json_preview(slides)}
    except Exception as exc:
        return {"kind": "presentation", "ok": False, "preview": f"preview_error={exc}"}


def _pdf_preview(path: Path) -> dict[str, Any]:
    try:
        import pdfplumber  # type: ignore

        with pdfplumber.open(str(path)) as pdf:
            text = "\n".join((page.extract_text() or "") for page in pdf.pages[:3])
        if text.strip():
            return {"kind": "pdf", "ok": True, "preview": _truncate(text)}
    except Exception:
        pass
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception:
        try:
            from PyPDF2 import PdfReader  # type: ignore
        except Exception:
            return {"kind": "pdf", "ok": False, "preview": "preview_unavailable=pypdf/PyPDF2 is not installed"}
    try:
        reader = PdfReader(str(path))
        text = "\n".join((page.extract_text() or "") for page in reader.pages[:3])
        return {"kind": "pdf", "ok": True, "preview": _truncate(text)}
    except Exception as exc:
        return {"kind": "pdf", "ok": False, "preview": f"preview_error={exc}"}


def _image_preview(path: Path) -> dict[str, Any]:
    try:
        from PIL import Image  # type: ignore
    except Exception:
        return {
            "kind": "image",
            "ok": True,
            "preview": "image file is available locally; inspect visually if needed.",
        }
    try:
        with Image.open(path) as image:
            info = {
                "format": image.format,
                "mode": image.mode,
                "width": image.width,
                "height": image.height,
                "animated": bool(getattr(image, "is_animated", False)),
            }
        return {
            "kind": "image",
            "ok": True,
            "preview": "image file is available locally; inspect visually if needed.",
            "metadata": info,
        }
    except Exception as exc:
        return {"kind": "image", "ok": False, "preview": f"preview_error={exc}"}


def _audio_preview(path: Path) -> dict[str, Any]:
    return {
        "kind": "audio",
        "ok": True,
        "preview": "audio file is available locally; not transcribed unless the user explicitly asks for audio-to-text.",
    }


def _zip_preview(path: Path) -> dict[str, Any]:
    try:
        with zipfile.ZipFile(path) as zf:
            members = []
            for info in zf.infolist()[:MAX_ARCHIVE_MEMBERS]:
                members.append({"name": info.filename, "size": info.file_size})
        return {"kind": "archive", "ok": True, "preview": _json_preview(members)}
    except Exception as exc:
        return {"kind": "archive", "ok": False, "preview": f"preview_error={exc}"}


def _seven_zip_preview(path: Path) -> dict[str, Any]:
    try:
        import py7zr  # type: ignore
    except Exception:
        return {"kind": "archive", "ok": False, "preview": "preview_unavailable=py7zr is not installed"}
    try:
        with py7zr.SevenZipFile(path, mode="r") as archive:
            names = archive.getnames()[:MAX_ARCHIVE_MEMBERS]
        return {"kind": "archive", "ok": True, "preview": _json_preview(names)}
    except Exception as exc:
        return {"kind": "archive", "ok": False, "preview": f"preview_error={exc}"}


def _odf_preview(path: Path) -> dict[str, Any]:
    try:
        from odf import teletype  # type: ignore
        from odf.opendocument import load  # type: ignore
        from odf.text import P  # type: ignore
    except Exception:
        return {"kind": "opendocument", "ok": False, "preview": "preview_unavailable=odfpy is not installed"}
    try:
        document = load(str(path))
        paragraphs = [teletype.extractText(p) for p in document.getElementsByType(P)]
        text = "\n".join(p for p in paragraphs if p.strip())
        return {"kind": "opendocument", "ok": True, "preview": _truncate(text)}
    except Exception as exc:
        return {"kind": "opendocument", "ok": False, "preview": f"preview_error={exc}"}


def analyze_path(path_like: str | Path) -> dict[str, Any]:
    path = Path(path_like).resolve()
    result: dict[str, Any] = {
        "path": str(path),
        "name": path.name,
        "suffix": path.suffix.lower(),
        "exists": path.exists(),
    }
    if not path.exists() or not path.is_file():
        result.update({"ok": False, "kind": "missing", "preview": "file does not exist"})
        return result

    stat = path.stat()
    result.update({"size": stat.st_size, "sha256": _sha256_file(path)})
    suffix = path.suffix.lower()
    if suffix == ".csv":
        detail = _csv_preview(path)
    elif suffix in {".txt", ".md", ".log", ".json", ".xml", ".yml", ".yaml", ".toml", ".ini", ".cfg", ".properties"}:
        detail = _read_text_preview(path)
    elif suffix in {".xlsx", ".xlsm"}:
        detail = _xlsx_preview(path)
    elif suffix == ".xls":
        detail = _xls_preview(path)
    elif suffix == ".docx":
        detail = _docx_preview(path)
    elif suffix == ".pptx":
        detail = _pptx_preview(path)
    elif suffix == ".pdf":
        detail = _pdf_preview(path)
    elif suffix in IMAGE_SUFFIXES:
        detail = _image_preview(path)
    elif suffix in AUDIO_SUFFIXES:
        detail = _audio_preview(path)
    elif suffix in {".zip", ".jar", ".mcpack", ".mcaddon"}:
        detail = _zip_preview(path)
    elif suffix == ".7z":
        detail = _seven_zip_preview(path)
    elif suffix in {".odt", ".ods", ".odp"}:
        detail = _odf_preview(path)
    else:
        detail = {"kind": "unknown", "ok": True, "preview": "local file exists but no previewer is configured for this extension."}
    result.update(detail)
    return result


def preview_path(path_like: str | Path) -> str:
    analysis = analyze_path(path_like)
    preview = str(analysis.get("preview") or "")
    metadata = analysis.get("metadata")
    if metadata:
        return preview + " metadata=" + json.dumps(metadata, ensure_ascii=False)
    return preview
