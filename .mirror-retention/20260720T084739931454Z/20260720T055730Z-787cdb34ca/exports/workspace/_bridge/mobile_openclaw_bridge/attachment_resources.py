"""Attachment resource handling for the mobile OpenClaw bridge.

Owns: attachment JSON parsing, safe attachment filenames, bounded attachment
materialization into the bridge attachment cache, file analysis previews, and
task attachment metadata helpers.
Non-goals: task permission policy, queue mutation, final reply sending, or
worker dispatch decisions.
State behavior: materialization may copy explicit attachments into the bridge
attachment cache and append resource-layer audit logs; preview helpers are
read-only.
Normal callers: mobile_openclaw_cli enqueue, prompt construction, supplement
handling, reply media preparation, and resource-layer smoke checks.
"""

from __future__ import annotations

import csv
import hashlib
import json
import re
import time
import unicodedata
import urllib.parse
import zipfile
from pathlib import Path
from typing import Any

from file_toolkit import analyze_path, preview_path
from resource_fetcher import (
    ResourceIntent,
    ResourceRequest,
    acquire_resource_with_policy,
    append_resource_log,
)

ROOT = Path(__file__).resolve().parent
ATTACHMENTS_DIR = ROOT / "attachments"
RESOURCE_LOG = ROOT / "logs" / "resource-fetcher.jsonl"
MAX_ATTACHMENT_PREVIEW_CHARS = 1200
MAX_ATTACHMENT_COPY_BYTES = 100 * 1024 * 1024

def parse_attachments_json(value: str | None) -> list[dict[str, Any]]:
    if not value:
        return []
    try:
        parsed = json.loads(value)
    except json.JSONDecodeError as exc:
        raise SystemExit(f"invalid --attachments-json: {exc}") from exc
    if isinstance(parsed, dict):
        parsed = [parsed]
    if not isinstance(parsed, list):
        raise SystemExit("--attachments-json must be a JSON array or object")
    result: list[dict[str, Any]] = []
    for item in parsed[:10]:
        if isinstance(item, dict):
            result.append(item)
        else:
            result.append({"value": str(item)})
    return result


def safe_filename(value: str, fallback: str = "attachment") -> str:
    name = Path(value or fallback).name.strip() or fallback
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "_", name)
    return name[:160] or fallback


def ascii_safe_filename(value: str, fallback: str = "attachment") -> str:
    raw = safe_filename(value, fallback=fallback)
    path = Path(raw)
    suffix = path.suffix
    stem = path.stem if suffix else raw
    ascii_stem = unicodedata.normalize("NFKD", stem).encode("ascii", "ignore").decode("ascii")
    ascii_stem = re.sub(r"[^A-Za-z0-9._-]+", "_", ascii_stem).strip("._-")
    if not ascii_stem:
        ascii_stem = fallback
    ascii_suffix = unicodedata.normalize("NFKD", suffix).encode("ascii", "ignore").decode("ascii")
    ascii_suffix = re.sub(r"[^A-Za-z0-9.]+", "", ascii_suffix)
    if suffix and not ascii_suffix.startswith("."):
        ascii_suffix = "." + ascii_suffix
    name = f"{ascii_stem[:96]}{ascii_suffix[:24]}"
    return name[:128] or fallback


def is_ascii_safe_filename(value: str) -> bool:
    name = Path(value or "").name
    if not name:
        return False
    try:
        name.encode("ascii")
    except UnicodeEncodeError:
        return False
    return name == safe_filename(name) and re.fullmatch(r"[A-Za-z0-9._ -]+", name) is not None


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as fh:
        for chunk in iter(lambda: fh.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def safe_local_path(value: Any) -> Path | None:
    if not isinstance(value, str) or not value.strip():
        return None
    try:
        path = Path(value.strip()).expanduser()
    except Exception:
        return None
    if not path.is_absolute():
        path = (ROOT / path).resolve()
    try:
        path = path.resolve()
    except Exception:
        return None
    if not path.exists() or not path.is_file():
        return None
    return path


def read_text_preview(path: Path, limit: int = MAX_ATTACHMENT_PREVIEW_CHARS) -> str:
    try:
        data = path.read_text(encoding="utf-8-sig", errors="replace")
    except Exception as exc:
        return f"preview_error={exc}"
    data = data.replace("\x00", "")
    return data[:limit] + ("...<truncated>" if len(data) > limit else "")


def csv_preview(path: Path, limit_rows: int = 8, limit_cols: int = 8) -> str:
    try:
        with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as fh:
            rows = []
            for index, row in enumerate(csv.reader(fh)):
                if index >= limit_rows:
                    break
                rows.append(row[:limit_cols])
    except Exception as exc:
        return f"preview_error={exc}"
    return json.dumps(rows, ensure_ascii=False)


def xlsx_preview(path: Path) -> str:
    try:
        import openpyxl  # type: ignore
    except Exception:
        return "preview_unavailable=openpyxl is not installed"
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for sheet in wb.worksheets[:3]:
            rows = []
            for index, row in enumerate(sheet.iter_rows(values_only=True)):
                if index >= 8:
                    break
                rows.append([cell for cell in row[:8]])
            parts.append({"sheet": sheet.title, "rows": rows})
        wb.close()
        return json.dumps(parts, ensure_ascii=False)[:MAX_ATTACHMENT_PREVIEW_CHARS]
    except Exception as exc:
        return f"preview_error={exc}"


def docx_preview(path: Path) -> str:
    try:
        import docx  # type: ignore
    except Exception:
        return "preview_unavailable=python-docx is not installed"
    try:
        doc = docx.Document(path)
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs[:30] if paragraph.text.strip())
        return text[:MAX_ATTACHMENT_PREVIEW_CHARS]
    except Exception as exc:
        return f"preview_error={exc}"


def pptx_preview(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore
    except Exception:
        return "preview_unavailable=python-pptx is not installed"
    try:
        prs = Presentation(path)
        slides = []
        for slide_index, slide in enumerate(prs.slides[:5], start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slides.append({"slide": slide_index, "text": texts[:8]})
        return json.dumps(slides, ensure_ascii=False)[:MAX_ATTACHMENT_PREVIEW_CHARS]
    except Exception as exc:
        return f"preview_error={exc}"


def pdf_preview(path: Path) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except Exception:
        try:
            from PyPDF2 import PdfReader  # type: ignore
        except Exception:
            return "preview_unavailable=pypdf/PyPDF2 is not installed"
    try:
        reader = PdfReader(str(path))
        text = "\n".join((page.extract_text() or "") for page in reader.pages[:3])
        return text[:MAX_ATTACHMENT_PREVIEW_CHARS]
    except Exception as exc:
        return f"preview_error={exc}"


def zip_listing_preview(path: Path) -> str:
    try:
        with zipfile.ZipFile(path) as zf:
            names = zf.namelist()[:30]
        return json.dumps(names, ensure_ascii=False)
    except Exception as exc:
        return f"preview_error={exc}"


def attachment_local_path(attachment: dict[str, Any]) -> Path | None:
    for key in ("local_path", "path", "file_path", "filePath", "temp_path", "tempPath"):
        path = safe_local_path(attachment.get(key))
        if path:
            return path
    return None


def attachment_url(attachment: dict[str, Any]) -> str:
    for key in ("url", "download_url", "downloadUrl", "remote_url", "remoteUrl"):
        value = str(attachment.get(key) or "").strip()
        if value:
            return value
    return ""


def attachment_display_name(attachment: dict[str, Any], fallback: str = "attachment") -> str:
    for key in ("name", "filename", "fileName", "original_name", "originalName"):
        value = str(attachment.get(key) or "").strip()
        if value:
            return value
    url = attachment_url(attachment)
    if url:
        parsed = urllib.parse.urlparse(url)
        name = Path(urllib.parse.unquote(parsed.path or "")).name
        if name:
            return name
    path = attachment_local_path(attachment)
    return path.name if path else fallback


def attachment_expected_sha256(attachment: dict[str, Any]) -> str:
    for key in ("expected_sha256", "checksum_sha256", "sha256_expected"):
        value = str(attachment.get(key) or "").strip().lower()
        if re.fullmatch(r"[0-9a-f]{64}", value):
            return value
    return ""


def persist_attachment_analysis(item: dict[str, Any], path: str) -> None:
    analysis = analyze_path(path)
    item["analysis_ok"] = bool(analysis.get("ok"))
    item["analysis_kind"] = analysis.get("kind")
    item["analysis_sha256"] = analysis.get("sha256")
    item["analysis_preview"] = str(analysis.get("preview") or "")[:MAX_ATTACHMENT_PREVIEW_CHARS]
    if analysis.get("metadata"):
        item["analysis_metadata"] = analysis.get("metadata")
    if analysis.get("error"):
        item["analysis_error"] = analysis.get("error")


def materialize_resource_attachment(attachment: dict[str, Any], index: int) -> dict[str, Any]:
    item = dict(attachment)
    source_path = attachment_local_path(item)
    url = attachment_url(item)
    original_name = attachment_display_name(item, fallback=f"attachment-{index}")
    request_metadata = {
        "attachment_index": index,
        "attachment_keys": sorted(str(key) for key in item.keys()),
    }
    result = None
    if source_path:
        result = acquire_resource_with_policy(
            ResourceRequest(
                source="mobile-openclaw-attachment",
                local_path=source_path,
                target_dir=ATTACHMENTS_DIR / time.strftime("%Y%m%d"),
                name=original_name,
                expected_sha256=attachment_expected_sha256(item),
                max_bytes=MAX_ATTACHMENT_COPY_BYTES,
                metadata=request_metadata,
            ),
            intent=ResourceIntent.EXPLICIT_ATTACHMENT,
        )
    elif url:
        result = acquire_resource_with_policy(
            ResourceRequest(
                source="mobile-openclaw-attachment-url",
                url=url,
                target_dir=ATTACHMENTS_DIR / time.strftime("%Y%m%d"),
                name=original_name,
                expected_sha256=attachment_expected_sha256(item),
                max_bytes=MAX_ATTACHMENT_COPY_BYTES,
                metadata=request_metadata,
            ),
            intent=ResourceIntent.EXPLICIT_ATTACHMENT,
        )
    else:
        item["resource_status"] = "skipped"
        item["resource_decision"] = "skipped"
        item["resource_policy"] = "explicit_attachment_v1"
        item["resource_policy_reason"] = "no_resource_reference"
        item["resource_next_action"] = "no_resource_materialization_needed"
        item["resource_error"] = "no local path or supported attachment URL"
        return item
    append_resource_log(RESOURCE_LOG, result)
    if result.ok:
        item["resource_status"] = "stored"
    elif result.decision == "deferred":
        item["resource_status"] = "deferred"
    else:
        item["resource_status"] = "failed"
    item["resource_source"] = result.resource_kind or ("local_file" if source_path else "url")
    item["resource_decision"] = result.decision
    item["resource_policy"] = result.policy_name
    item["resource_policy_reason"] = result.policy_reason
    item["resource_next_action"] = result.next_action
    item["resource_risk_flags"] = list(result.risk_flags or [])
    item["resource_cache_hit"] = result.cache_hit
    item["resource_error"] = result.error
    item["resource_metadata"] = result.metadata or {}
    if result.original_local_path:
        item["original_local_path"] = result.original_local_path
    if result.local_path:
        item["local_path"] = result.local_path
    if result.stored_path:
        item["stored_path"] = result.stored_path
    if result.sha256:
        item["sha256"] = result.sha256
    if result.size:
        item["size"] = result.size
    if result.ok and result.stored_path:
        item["stored"] = True
        persist_attachment_analysis(item, result.stored_path)
    elif source_path:
        try:
            item.setdefault("size", source_path.stat().st_size)
        except OSError:
            pass
        item["storage_note"] = f"not copied; {result.error}"
    return item


def materialize_attachments(attachments: list[dict[str, Any]]) -> list[dict[str, Any]]:
    materialized: list[dict[str, Any]] = []
    for index, attachment in enumerate(attachments[:10], start=1):
        try:
            item = materialize_resource_attachment(attachment, index)
        except Exception as exc:
            item = dict(attachment)
            item["resource_status"] = "failed"
            item["resource_error"] = f"materialize_failed={exc}"
            item["storage_note"] = f"copy_failed={exc}"
        materialized.append(item)
    return materialized


def describe_attachment(attachment: dict[str, Any], index: int) -> list[str]:
    name = str(attachment.get("name") or attachment.get("filename") or attachment.get("fileName") or "")
    kind = str(attachment.get("type") or attachment.get("kind") or attachment.get("mediaType") or "")
    mime = str(attachment.get("mime") or attachment.get("mimeType") or "")
    size = attachment.get("size") or attachment.get("fileSize") or ""
    path = attachment_local_path(attachment)
    url = str(attachment.get("url") or attachment.get("download_url") or attachment.get("downloadUrl") or "")

    lines = [f"  attachment[{index}]: type={kind or 'unknown'} name={name or '<unnamed>'} mime={mime or '<unknown>'} size={size or '<unknown>'}"]
    resource_status = str(attachment.get("resource_status") or "").strip()
    if resource_status:
        lines.append(
            "    resource_status={status} source={source} decision={decision} policy={policy} cache_hit={cache_hit} error={error}".format(
                status=resource_status,
                source=attachment.get("resource_source") or "-",
                decision=attachment.get("resource_decision") or "-",
                policy=attachment.get("resource_policy") or "-",
                cache_hit=attachment.get("resource_cache_hit"),
                error=attachment.get("resource_error") or "-",
            )
        )
    if path:
        lines.append(f"    local_path={path}")
        analysis_kind = attachment.get("analysis_kind")
        analysis_sha = attachment.get("analysis_sha256") or attachment.get("sha256")
        analysis_preview = str(attachment.get("analysis_preview") or "")
        if not analysis_kind or not analysis_preview:
            analysis = analyze_path(path)
            analysis_kind = analysis.get("kind")
            analysis_sha = analysis.get("sha256") or analysis_sha
            analysis_preview = str(analysis.get("preview") or preview_path(path))
        lines.append(f"    analysis_kind={analysis_kind} sha256={str(analysis_sha or '')[:16]}")
        lines.append("    preview=" + analysis_preview[:MAX_ATTACHMENT_PREVIEW_CHARS])
    elif url:
        lines.append(f"    remote_url_present=true url={url} preview=resource_not_available")
    else:
        keys = ",".join(sorted(str(key) for key in attachment.keys())[:20])
        lines.append(f"    preview=no local path or URL available; metadata_keys={keys}")
    return lines


def task_attachments(task: dict[str, Any]) -> list[dict[str, Any]]:
    raw = task.get("attachments_json")
    if not raw:
        return []
    try:
        parsed = json.loads(str(raw))
    except json.JSONDecodeError:
        return [{"type": "parse_error", "name": "attachments_json", "value": str(raw)[:200]}]
    if isinstance(parsed, dict):
        parsed = [parsed]
    if not isinstance(parsed, list):
        return []
    return [item for item in parsed if isinstance(item, dict)]


def task_has_attachments(task: dict[str, Any]) -> bool:
    if task_attachments(task):
        return True
    raw_metadata = task.get("metadata_json")
    if not raw_metadata:
        return False
    try:
        metadata = json.loads(str(raw_metadata))
    except json.JSONDecodeError:
        return False
    if not isinstance(metadata, dict):
        return False
    try:
        return int(metadata.get("attachment_count") or 0) > 0
    except (TypeError, ValueError):
        return False
